"""
Multi-format document loader.
Extracts text from PDF, DOCX, TXT, PPTX, and CSV files
while preserving page/slide numbers for citation tracking.
"""

import csv

from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional


@dataclass
class ExtractedPage:
    """A single page/slide/section of extracted text."""
    text: str
    page_number: int
    document_name: str
    metadata: dict = field(default_factory=dict)


def load_document(file_path: str, document_name: Optional[str] = None) -> list[ExtractedPage]:
    """
    Load a document and extract text with page-level granularity.

    Args:
        file_path: Path to the document file.
        document_name: Display name for the document (defaults to filename).

    Returns:
        List of ExtractedPage objects with text and metadata.

    Raises:
        ValueError: If the file format is not supported.
        FileNotFoundError: If the file does not exist.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    if document_name is None:
        document_name = path.name

    extension = path.suffix.lower()

    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".txt": _extract_txt,
        ".pptx": _extract_pptx,
        ".csv": _extract_csv,
        ".xlsx": _extract_xlsx,
    }

    extractor = extractors.get(extension)
    if extractor is None:
        raise ValueError(
            f"Unsupported file format: {extension}. "
            f"Supported formats: {', '.join(extractors.keys())}"
        )

    pages = extractor(str(path), document_name)

    # Filter out empty pages
    pages = [p for p in pages if p.text.strip()]

    return pages


def _extract_pdf(file_path: str, document_name: str) -> list[ExtractedPage]:
    """Extract text from PDF using PyMuPDF, one ExtractedPage per PDF page."""
    import fitz  # PyMuPDF

    pages = []
    doc = fitz.open(file_path)

    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")

            if text.strip():
                pages.append(ExtractedPage(
                    text=text.strip(),
                    page_number=page_num + 1,  # 1-indexed
                    document_name=document_name,
                    metadata={
                        "format": "pdf",
                        "total_pages": len(doc),
                        "width": page.rect.width,
                        "height": page.rect.height,
                    }
                ))
    finally:
        doc.close()

    return pages


def _extract_docx(file_path: str, document_name: str) -> list[ExtractedPage]:
    """
    Extract text from DOCX files.
    Groups paragraphs into approximate pages (~3000 chars per page).
    """
    from docx import Document

    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    if not paragraphs:
        return []

    # Group paragraphs into approximate pages
    pages = []
    current_text = []
    current_length = 0
    page_num = 1
    chars_per_page = 3000  # approximate page length

    for para in paragraphs:
        current_text.append(para)
        current_length += len(para)

        if current_length >= chars_per_page:
            pages.append(ExtractedPage(
                text="\n\n".join(current_text),
                page_number=page_num,
                document_name=document_name,
                metadata={"format": "docx"}
            ))
            current_text = []
            current_length = 0
            page_num += 1

    # Don't forget remaining text
    if current_text:
        pages.append(ExtractedPage(
            text="\n\n".join(current_text),
            page_number=page_num,
            document_name=document_name,
            metadata={"format": "docx"}
        ))

    # Update total pages in metadata
    total = len(pages)
    for p in pages:
        p.metadata["total_pages"] = total

    return pages


def _extract_txt(file_path: str, document_name: str) -> list[ExtractedPage]:
    """
    Extract text from plain text files.
    Groups lines into approximate pages (~3000 chars per page).
    """
    # Try multiple encodings
    text = None
    for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                text = f.read()
            break
        except (UnicodeDecodeError, UnicodeError):
            continue

    if text is None:
        raise ValueError(f"Could not decode text file: {file_path}")

    if not text.strip():
        return []

    # Split into approximate pages
    pages = []
    chars_per_page = 3000
    page_num = 1

    for i in range(0, len(text), chars_per_page):
        chunk = text[i:i + chars_per_page].strip()
        if chunk:
            pages.append(ExtractedPage(
                text=chunk,
                page_number=page_num,
                document_name=document_name,
                metadata={"format": "txt", "total_pages": 0}
            ))
            page_num += 1

    # Update total pages
    total = len(pages)
    for p in pages:
        p.metadata["total_pages"] = total

    return pages


def _extract_pptx(file_path: str, document_name: str) -> list[ExtractedPage]:
    """Extract text from PowerPoint files, one ExtractedPage per slide."""
    from pptx import Presentation

    prs = Presentation(file_path)
    pages = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        texts.append(para_text)

            # Also extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))

        if texts:
            pages.append(ExtractedPage(
                text="\n".join(texts),
                page_number=slide_num,
                document_name=document_name,
                metadata={
                    "format": "pptx",
                    "total_pages": len(prs.slides),
                    "slide_title": texts[0][:100] if texts else "",
                }
            ))

    return pages


def _extract_csv(file_path: str, document_name: str) -> list[ExtractedPage]:
    """
    Extract text from CSV files.
    Groups rows into pages (~50 rows per page) with header context.
    """
    rows = []

    # Try multiple encodings
    for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
        try:
            with open(file_path, "r", encoding=encoding, newline="") as f:
                reader = csv.reader(f)
                rows = list(reader)
            break
        except (UnicodeDecodeError, UnicodeError):
            continue

    if not rows:
        return []

    header = rows[0] if rows else []
    data_rows = rows[1:] if len(rows) > 1 else rows

    pages = []
    rows_per_page = 50
    page_num = 1

    for i in range(0, len(data_rows), rows_per_page):
        batch = data_rows[i:i + rows_per_page]

        # Format rows with header context
        text_lines = []
        if header:
            text_lines.append("Columns: " + " | ".join(header))
            text_lines.append("-" * 40)

        for row in batch:
            if header and len(row) == len(header):
                row_text = "; ".join(
                    f"{h}: {v}" for h, v in zip(header, row) if v.strip()
                )
            else:
                row_text = " | ".join(v for v in row if v.strip())

            if row_text:
                text_lines.append(row_text)

        if text_lines:
            pages.append(ExtractedPage(
                text="\n".join(text_lines),
                page_number=page_num,
                document_name=document_name,
                metadata={
                    "format": "csv",
                    "total_rows": len(data_rows),
                    "columns": header,
                }
            ))
            page_num += 1

    # Update total pages
    total = len(pages)
    for p in pages:
        p.metadata["total_pages"] = total

    return pages


def _extract_xlsx(file_path: str, document_name: str) -> list[ExtractedPage]:
    """Extract text from Excel XLSX files."""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    pages = []
    page_num = 1

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        
        # Group rows into approximate pages
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
            
        header = [str(c) for c in rows[0] if c is not None] if rows else []
        data_rows = rows[1:] if len(rows) > 1 else rows
        
        rows_per_page = 50
        
        for i in range(0, len(data_rows), rows_per_page):
            batch = data_rows[i:i + rows_per_page]
            text_lines = [f"Sheet: {sheet_name}"]
            
            if header:
                text_lines.append("Columns: " + " | ".join(header))
                text_lines.append("-" * 40)
                
            for row in batch:
                row_text = " | ".join(str(v) for v in row if v is not None and str(v).strip())
                if row_text:
                    text_lines.append(row_text)
                    
            if len(text_lines) > (2 if header else 1):
                pages.append(ExtractedPage(
                    text="\n".join(text_lines),
                    page_number=page_num,
                    document_name=document_name,
                    metadata={
                        "format": "xlsx",
                        "sheet": sheet_name,
                    }
                ))
                page_num += 1

    total = len(pages)
    for p in pages:
        p.metadata["total_pages"] = total

    return pages
